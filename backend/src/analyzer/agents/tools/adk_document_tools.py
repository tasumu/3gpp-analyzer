"""Document tools for ADK-based meeting analysis agents."""

import asyncio
import io
import logging
from typing import Any

from google.adk.tools import ToolContext

from analyzer.agents.context import AgentToolContext, get_current_agent_context

logger = logging.getLogger(__name__)


def _extract_docx_text(content: bytes) -> str:
    """Extract text from .docx bytes using python-docx.

    Returns paragraphs and tables formatted as markdown.
    """
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(content))
    parts: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    for table in doc.tables:
        rows_text: list[str] = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            rows_text.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows_text.append("| " + " | ".join(["---"] * len(cells)) + " |")
        if rows_text:
            parts.append("\n".join(rows_text))

    return "\n\n".join(parts)


def _extract_pptx_text(content: bytes) -> str:
    """Extract text from .pptx bytes using python-pptx.

    Returns slide content formatted as markdown with slide numbers as headings.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                rows_text: list[str] = []
                for j, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_text.append("| " + " | ".join(cells) + " |")
                    if j == 0:
                        rows_text.append("| " + " | ".join(["---"] * len(cells)) + " |")
                if rows_text:
                    slide_texts.append("\n".join(rows_text))
        if slide_texts:
            parts.append(f"## Slide {i}\n\n" + "\n\n".join(slide_texts))

    return "\n\n".join(parts)


def _extract_xlsx_text(content: bytes) -> str:
    """Extract text from .xlsx bytes using openpyxl.

    Returns sheet content formatted as markdown tables.
    """
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sections: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        section_parts: list[str] = [f"## Sheet: {sheet_name}\n"]

        header = rows[0]
        header_strs = [str(c) if c is not None else "" for c in header]
        section_parts.append("| " + " | ".join(header_strs) + " |")
        section_parts.append("| " + " | ".join(["---"] * len(header)) + " |")

        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            section_parts.append("| " + " | ".join(cells) + " |")

        sections.append("\n".join(section_parts))

    wb.close()
    return "\n\n".join(sections)


async def get_document_summary(
    document_id: str,
    tool_context: ToolContext = None,
) -> dict[str, Any]:
    """
    Get the summary of a specific document.

    Returns the pre-computed analysis summary if available,
    or basic document metadata if not analyzed yet.

    Args:
        document_id: The document ID to get summary for.
            This is typically the contribution number.
        tool_context: ADK tool context (injected automatically by ADK).

    Returns:
        Document summary and metadata including contribution_number, title,
        source, status, and summary text.
    """
    # Get context from contextvar (preferred) or ADK's state (fallback)
    ctx: AgentToolContext | None = get_current_agent_context()
    if not ctx and tool_context and tool_context.state:
        ctx = tool_context.state.get("agent_context")

    if not ctx or not ctx.firestore:
        return {"error": "Firestore not available"}

    logger.info(f"Getting summary for document: {document_id}")

    try:
        # Get document metadata
        doc_data = await ctx.firestore.get_document(document_id)
        if not doc_data:
            return {"error": f"Document not found: {document_id}"}

        result = {
            "document_id": document_id,
            "contribution_number": doc_data.get("contribution_number", ""),
            "title": doc_data.get("title", "Untitled"),
            "source": doc_data.get("source", "Unknown"),
            "status": doc_data.get("status", "unknown"),
        }

        # Try to get cached summary from document_summaries collection
        try:
            language = ctx.language if ctx else "ja"
            cache_key = f"{document_id}_{language}"
            doc_ref = ctx.firestore.client.collection("document_summaries").document(cache_key)
            doc = await asyncio.to_thread(doc_ref.get)
            if doc.exists:
                data = doc.to_dict()
                result["summary"] = data.get("summary", "No summary available")
                result["key_points"] = data.get("key_points", [])
                result["has_analysis"] = True
            else:
                result["summary"] = "No analysis available for this document"
                result["has_analysis"] = False
        except Exception as e:
            logger.warning(f"Error fetching summary for {document_id}: {e}")
            result["summary"] = "Unable to retrieve analysis"
            result["has_analysis"] = False

        return result

    except Exception as e:
        logger.error(f"Error getting summary for document {document_id}: {e}")
        return {"error": str(e)}


async def get_document_content(
    document_id: str,
    max_chunks: int = 500,
    tool_context: ToolContext = None,
) -> dict[str, Any]:
    """
    Get the full content of a specific document.

    For indexed documents, returns all chunks organized by sections.
    For non-indexed documents, falls back to direct text extraction from GCS.
    Supports .docx, .doc (via normalized .docx), .xlsx, and .pptx.

    Args:
        document_id: The document ID to get content for.
        max_chunks: Maximum number of chunks to return. Default: 500.
        tool_context: ADK tool context (injected automatically by ADK).

    Returns:
        Document content organized by sections with clause numbers,
        titles, content text, and page numbers.
    """
    # Get context from contextvar (preferred) or ADK's state (fallback)
    ctx: AgentToolContext | None = get_current_agent_context()
    if not ctx and tool_context and tool_context.state:
        ctx = tool_context.state.get("agent_context")

    if not ctx:
        return {"error": "Agent context not initialized", "sections": [], "total_chunks": 0}

    logger.info(f"Getting content for document: {document_id}")

    try:
        evidences = await ctx.evidence_provider.get_by_document(
            document_id=document_id,
            top_k=max_chunks,
        )

        # Fallback: if no chunks exist, try reading the original file from GCS
        if not evidences and ctx.storage and ctx.firestore:
            return await _get_document_content_from_gcs(ctx, document_id)

        # Organize by clause
        content_sections = []
        for ev in evidences:
            content_sections.append(
                {
                    "clause": ev.clause_number or "Unknown",
                    "title": ev.clause_title or "",
                    "content": ev.content,
                    "page": ev.page_number,
                }
            )

        # Track these as used evidences
        ctx.used_evidences.extend(evidences)

        return {
            "document_id": document_id,
            "sections": content_sections,
            "total_chunks": len(content_sections),
        }

    except Exception as e:
        logger.error(f"Error getting content for document {document_id}: {e}")
        return {"error": str(e), "sections": [], "total_chunks": 0}


async def _get_document_content_from_gcs(ctx: AgentToolContext, document_id: str) -> dict[str, Any]:
    """Fallback: read document content directly from GCS for non-indexed documents."""
    doc_data = await ctx.firestore.get_document(document_id)
    if not doc_data:
        return {
            "error": f"Document not found: {document_id}",
            "sections": [],
            "total_chunks": 0,
        }

    source_file = doc_data.get("source_file", {})
    gcs_path = source_file.get("gcs_normalized_path") or source_file.get("gcs_original_path")
    filename = source_file.get("filename", "")

    if not gcs_path:
        return {
            "error": "Document file not available in GCS (not yet downloaded)",
            "sections": [],
            "total_chunks": 0,
        }

    # Determine extraction strategy based on actual file at gcs_path and original filename
    gcs_ext = gcs_path.lower().rsplit(".", 1)[-1] if "." in gcs_path else ""
    filename_ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    extractor = None
    if gcs_ext == "docx":
        # Covers both .docx originals and .doc files normalized to .docx
        extractor = _extract_docx_text
    elif filename_ext == "xlsx":
        extractor = _extract_xlsx_text
    elif filename_ext == "pptx":
        extractor = _extract_pptx_text

    if extractor is None:
        return {
            "error": (
                f"Direct reading not supported for {filename} "
                f"(supported: .docx, .doc, .xlsx, .pptx)"
            ),
            "sections": [],
            "total_chunks": 0,
        }

    logger.info(f"Reading non-indexed document from GCS: {gcs_path} (original: {filename})")
    content_bytes = await ctx.storage.download_bytes(gcs_path)
    text = extractor(content_bytes)

    return {
        "document_id": document_id,
        "sections": [{"clause": "Full Document", "title": "", "content": text, "page": None}],
        "total_chunks": 1,
        "source": "gcs_fallback",
    }
